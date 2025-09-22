#!/usr/bin/env python3
"""
VirtuLearn Principal Dashboard - Main Application File
Clean UI-focused implementation that uses abstracted evaluation models.
"""

import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from datetime import datetime, date
from dotenv import load_dotenv
import os
import json
import hashlib
import openai
import io
from pathlib import Path

# Try to import optional file parsing libraries
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Import evaluation functions from model package
from model import (
    run_evaluation_sync,
    generate_comprehensive_evaluation_report
)

# Import data manager
from utils.data_manager import LectureDataManager

# Load environment variables
load_dotenv()

openai.api_key = os.getenv('OPENAI_API_KEY')

# Initialize data manager
@st.cache_resource
def get_data_manager():
    return LectureDataManager(use_mongodb=True)

# Page configuration
st.set_page_config(
    page_title="VirtuLearn - Principal Lecture Evaluation System",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded",
    menu_items={
        'Get Help': 'https://www.streamlit.io/community',
        'Report a bug': "https://github.com/yourusername/virtulearn/issues",
        'About': "# VirtuLearn Principal Dashboard\nLecture Quality Evaluation & Analytics Platform"
    }
)

# Custom CSS for better styling
st.markdown("""
<style>
    .main-header {
        font-size: 3rem;
        color: #1f77b4;
        text-align: center;
        margin-bottom: 2rem;
    }
    .metric-container {
        background-color: #f0f2f6;
        padding: 1rem;
        border-radius: 10px;
        margin: 0.5rem 0;
    }
    .score-excellent {
        background: linear-gradient(90deg, #4CAF50, #45a049);
        color: white;
        padding: 1rem;
        border-radius: 10px;
        text-align: center;
    }
    .score-good {
        background: linear-gradient(90deg, #2196F3, #1976D2);
        color: white;
        padding: 1rem;
        border-radius: 10px;
        text-align: center;
    }
    .score-needs-improvement {
        background: linear-gradient(90deg, #FF9800, #F57C00);
        color: white;
        padding: 1rem;
        border-radius: 10px;
        text-align: center;
    }
    .upload-section {
        background-color: #f8f9fa;
        padding: 1.5rem;
        border-radius: 10px;
        border-left: 4px solid #1f77b4;
        margin-bottom: 1rem;
    }
</style>
""", unsafe_allow_html=True)

# File parsing utility functions
def parse_text_file(file):
    """Parse text file content"""
    try:
        content = str(file.read(), "utf-8")
        return content
    except Exception as e:
        st.error(f"Error reading text file: {str(e)}")
        return ""

def parse_pdf_file(file):
    """Parse PDF file content using PyMuPDF"""
    if not PDF_AVAILABLE:
        st.warning("PDF parsing not available. Please install PyMuPDF: pip install PyMuPDF")
        return ""
    
    try:
        # Read the file content
        file_bytes = file.read()
        
        # Open PDF document from bytes
        pdf_document = fitz.open(stream=file_bytes, filetype="pdf")
        
        text = ""
        for page_num in range(pdf_document.page_count):
            page = pdf_document[page_num]
            text += page.get_text() + "\n"
        
        pdf_document.close()
        return text
    except Exception as e:
        st.error(f"Error reading PDF file: {str(e)}")
        return ""

def parse_docx_file(file):
    """Parse DOCX file content"""
    if not DOCX_AVAILABLE:
        st.warning("DOCX parsing not available. Please install python-docx: pip install python-docx")
        return ""
    
    try:
        doc = Document(io.BytesIO(file.read()))
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        st.error(f"Error reading DOCX file: {str(e)}")
        return ""

def parse_pptx_file(file):
    """Parse PPTX file content"""
    if not PPTX_AVAILABLE:
        st.warning("PPTX parsing not available. Please install python-pptx: pip install python-pptx")
        return ""
    
    try:
        prs = Presentation(io.BytesIO(file.read()))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"Error reading PPTX file: {str(e)}")
        return ""

def parse_uploaded_file(uploaded_file):
    """Parse uploaded file based on file type"""
    if uploaded_file is None:
        return ""
    
    file_type = uploaded_file.type
    
    if file_type == "text/plain":
        return parse_text_file(uploaded_file)
    elif file_type == "application/pdf":
        return parse_pdf_file(uploaded_file)
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return parse_docx_file(uploaded_file)
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return parse_pptx_file(uploaded_file)
    elif file_type == "application/vnd.ms-powerpoint":
        st.warning("Please convert .ppt files to .pptx format for parsing")
        return ""
    else:
        st.warning(f"Unsupported file type: {file_type}")
        return ""

def show_lecture_upload():
    """Display lecture upload and evaluation interface for principal"""
    st.header("📚 Lecture Evaluation System")
    
    st.markdown("""
    <div class="upload-section">
        <h3>🎯 Upload Lecture Materials for Quality Assessment</h3>
        <p>Upload teacher lecture transcripts and materials to evaluate teaching quality and engagement</p>
    </div>
    """, unsafe_allow_html=True)
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.subheader("Lecture Information")
        teacher_name = st.text_input("Teacher Name:", placeholder="e.g., Dr. Smith")
        lecture_title = st.text_input("Lecture Title:", placeholder="e.g., Introduction to Calculus")
        lecture_date = st.date_input("Lecture Date:", value=datetime.now().date())
        course_code = st.text_input("Course Code:", placeholder="e.g., MATH101")
        duration = st.number_input("Duration (minutes):", min_value=1, max_value=180, value=50)
        class_size = st.number_input("Class Size:", min_value=1, max_value=200, value=25)
        
        st.subheader("Expected Content")
        topics_covered = st.text_area("Topics to be Covered:", placeholder="List expected topics separated by commas")
        learning_objectives = st.text_area("Learning Objectives:", placeholder="What should students learn?")
    
    with col2:
        st.subheader("Upload Materials")
        
        # Transcript upload (required)
        transcript_file = st.file_uploader(
            "📝 Upload Lecture Transcript (Required)",
            type=['txt', 'docx', 'pdf'],
            help="Upload the lecture transcript for evaluation"
        )
        
        # Slides upload (optional)
        slides_file = st.file_uploader(
            "🖼️ Upload Lecture Slides (Optional)",
            type=['pptx', 'pdf', 'ppt'],
            help="Upload lecture slides for additional context"
        )
        
        # Source materials (required)
        materials_files = st.file_uploader(
            "📎 Upload Source Materials (Required)",
            type=['pdf', 'docx', 'xlsx', 'txt'],
            accept_multiple_files=True,
            help="Upload handouts or reference materials"
        )
    
    if st.button("🔍 Evaluate Lecture Quality", type="primary"):
        if transcript_file and teacher_name and lecture_title:
            with st.spinner("Evaluating lecture quality using AI agents..."):
                try:
                    # Get data manager
                    data_manager = get_data_manager()
                    
                    # Process transcript file
                    transcript_text = parse_uploaded_file(transcript_file)
                    if not transcript_text:
                        st.error("Failed to parse transcript file. Please try a different format.")
                        return
                    
                    # Process slides file if uploaded
                    slides_content = ""
                    if slides_file:
                        slides_content = parse_uploaded_file(slides_file)
                        if not slides_content:
                            st.warning("Could not parse slides file, continuing without slides content.")
                    
                    # Process source materials if uploaded
                    source_materials_content = ""
                    if materials_files:
                        st.write(f"🔍 DEBUG: Processing {len(materials_files)} material files")
                        for material_file in materials_files:
                            st.write(f"🔍 DEBUG: Processing file: {material_file.name}, type: {material_file.type}, size: {material_file.size}")
                            material_content = parse_uploaded_file(material_file)
                            if material_content:
                                source_materials_content += f"\n\n--- {material_file.name} ---\n"
                                source_materials_content += material_content
                                st.write(f"✅ Parsed {material_file.name}: {len(material_content)} characters")
                            else:
                                st.warning(f"Could not parse {material_file.name}, skipping this file.")
                    else:
                        st.info("ℹ️ No source materials uploaded. Upload textbooks, papers, or reference materials for AI-powered fact checking.")
                    
                    st.write(f"🔍 DEBUG: Total source materials: {len(source_materials_content)} characters")
                    if source_materials_content:
                        st.write(f"🔍 DEBUG: Source materials preview: {source_materials_content[:200]}...")
                    
                    # Calculate lecture quality score using abstracted evaluation system
                    try:
                        with st.spinner('Analyzing lecture... This may take a few minutes.'):
                            score, score_components, analysis_details = run_evaluation_sync(
                                transcript_text, 
                                topics_covered, 
                                duration,
                                source_materials_content,
                                slides_content
                            )
                    except Exception as e:
                        st.error(f"Error during evaluation: {str(e)}")
                        st.warning("Falling back to basic evaluation...")
                        
                        # Fallback to basic evaluation
                        try:
                            from model.lecture_evaluator import run_fallback_evaluation
                            score, score_components, analysis_details = run_fallback_evaluation(
                                transcript_text, 
                                topics_covered, 
                                duration,
                                source_materials_content,
                                slides_content
                            )
                            st.success("✅ Basic lecture analysis completed!")
                        except Exception as fallback_error:
                            st.error(f"Error during fallback evaluation: {str(fallback_error)}")
                            return
                    
                    # Generate comprehensive evaluation report
                    evaluation_report = generate_comprehensive_evaluation_report(
                        score, score_components, transcript_text, topics_covered, analysis_details
                    )
                    
                    # Create lecture entry in database
                    lecture_entry_data = {
                        'title': lecture_title,
                        'teacher_id': teacher_name,
                        'course_code': course_code,
                        'date': lecture_date,
                        'transcript_text': transcript_text,
                        'duration': int(duration) if duration else None,
                        'topics': [topic.strip() for topic in topics_covered.split(",")] if topics_covered else [],
                        'objectives': [obj.strip() for obj in learning_objectives.split(",")] if learning_objectives else []
                    }
                    
                    # Add parsed content if available
                    if source_materials_content:
                        lecture_entry_data['source_materials'] = source_materials_content
                    if slides_content:
                        lecture_entry_data['slides_content'] = slides_content
                    
                    lecture_id = data_manager.create_lecture_entry(**lecture_entry_data)
                    
                    # Store evaluation results in database
                    evaluation_data = {
                        'lecture_id': lecture_id,
                        'teacher_name': teacher_name,
                        'course_code': course_code,
                        'lecture_title': lecture_title,
                        'evaluation_score': score,
                        'score_breakdown': score_components,
                        'evaluation_report': evaluation_report,
                        'analysis_details': analysis_details,
                        'class_size': class_size,
                        'evaluation_timestamp': datetime.now(),
                        'evaluated_by': 'Principal'
                    }
                    
                    # Store evaluation in database
                    if hasattr(data_manager, 'store_evaluation'):
                        eval_id = data_manager.store_evaluation(evaluation_data)
                    else:
                        # Store as metadata if store_evaluation doesn't exist
                        eval_id = data_manager.save_lecture_data(lecture_id, evaluation_data, 'evaluation')
                    
                    # Store additional files if uploaded (need to re-read since we already parsed them)
                    if slides_file:
                        # Reset file pointer and read again for storage
                        slides_file.seek(0)
                        file_id = data_manager.store_uploaded_file(
                            lecture_id=lecture_id,
                            file_content=slides_file.read(),
                            filename=slides_file.name,
                            file_type="slides"
                        )
                    
                    if materials_files:
                        for material_file in materials_files:
                            # Reset file pointer and read again for storage
                            material_file.seek(0)
                            file_id = data_manager.store_uploaded_file(
                                lecture_id=lecture_id,
                                file_content=material_file.read(),
                                filename=material_file.name,
                                file_type="material"
                            )
                    
                    # Display evaluation results
                    st.success(f"✅ Lecture evaluation completed! Evaluation ID: {eval_id}")
                    display_evaluation_results(score, score_components, evaluation_report, analysis_details, 
                                             teacher_name, course_code, lecture_date, duration, class_size)
                    
                except Exception as e:
                    st.error(f"❌ Error during evaluation: {str(e)}")
                    st.error("Please check your files and try again.")
        else:
            st.warning("⚠️ Please provide teacher name, lecture title, and transcript file to proceed.")

def display_evaluation_results(score, score_components, evaluation_report, analysis_details, 
                             teacher_name, course_code, lecture_date, duration, class_size):
    """Display the evaluation results in a structured format"""
    
    # Show evaluation score with color coding
    st.subheader("🎯 Lecture Quality Evaluation")
    
    score_col1, score_col2, score_col3 = st.columns([2, 1, 1])
    
    with score_col1:
        if score >= 85:
            st.markdown(f"""
            <div class="score-excellent">
                <h2>Overall Score: {score:.1f}/100</h2>
                <h3>Grade: {evaluation_report['grade']} - Excellent!</h3>
                <p>🌟 Outstanding lecture quality with AI-verified standards</p>
            </div>
            """, unsafe_allow_html=True)
        elif score >= 70:
            st.markdown(f"""
            <div class="score-good">
                <h2>Overall Score: {score:.1f}/100</h2>
                <h3>Grade: {evaluation_report['grade']} - Good</h3>
                <p>👍 Solid lecture with AI-identified improvement areas</p>
            </div>
            """, unsafe_allow_html=True)
        else:
            st.markdown(f"""
            <div class="score-needs-improvement">
                <h2>Overall Score: {score:.1f}/100</h2>
                <h3>Grade: {evaluation_report['grade']} - Needs Improvement</h3>
                <p>⚠️ Requires significant enhancement per AI analysis</p>
            </div>
            """, unsafe_allow_html=True)
    
    with score_col2:
        st.metric("Teacher", teacher_name, f"Class Size: {class_size}")
        st.metric("Duration", f"{duration} min", f"Words: {evaluation_report['word_count']:,}")
    
    with score_col3:
        st.metric("Course", course_code, f"Date: {lecture_date}")
        st.metric("Topics", f"{len(evaluation_report['topics_covered'])}", "Covered")
    
    # Score breakdown
    st.subheader("📊 Detailed Score Breakdown")
    
    # Define maximum possible scores for each component (corrected values)
    max_scores = {
        'Topic Coverage': 30,  # Topic Coverage
        'Engagement': 20,       
        'Correctness': 30       
    }
    
    # Calculate independent percentages based on each component's maximum
    breakdown_data = []
    for component, score in score_components.items():
        if component in max_scores:  # Only include components we want to keep
            max_score = max_scores[component]
            percentage = (score / max_score) * 100 if max_score > 0 else 0
            breakdown_data.append({
                'Component': component,
                'Score': f"{score:.1f}",
                'Max Score': max_score,
                'Percentage': f"{percentage:.1f}%"
            })
    
    breakdown_df = pd.DataFrame(breakdown_data)
    
    # Display table showing individual scores (Requirement 1)
    st.markdown("**1. Individual Scores Table**")
    st.dataframe(breakdown_df, hide_index=True, use_container_width=True)
    
    # Pie chart for correctness evaluation (Requirement 2)
    st.markdown("**2. Correctness Analysis - Claims Distribution**")
    correctness = analysis_details.get('correctness', {})
    scoring = correctness.get('scoring_details', {})
    
    if scoring and scoring.get('total_claims', 0) > 0:
        # Create pie chart for correct/incorrect/unfounded statements
        correct_claims = scoring.get('correct_claims', 0)
        incorrect_claims = scoring.get('incorrect_claims', 0)
        unsupported_claims = scoring.get('unsupported_claims', 0)
        
        claims_fig = go.Figure(data=[go.Pie(
            labels=['Correct', 'Incorrect', 'Unfounded'],
            values=[correct_claims, incorrect_claims, unsupported_claims],
            hole=0.3,
            marker_colors=['#4CAF50', '#F44336', '#FF9800']
        )])
        claims_fig.update_layout(
            title=f"Claims Analysis ({correct_claims + incorrect_claims + unsupported_claims} total claims)",
            height=400,
            showlegend=True
        )
        st.plotly_chart(claims_fig, use_container_width=True)
        
        # Show detailed claims analysis in an expandable section
        with st.expander("🔍 View Detailed Claims Analysis"):
            claims_list = correctness.get('claims', [])
            if claims_list:
                for i, claim in enumerate(claims_list):
                    claim_text = claim.get('claim', 'Unknown claim')
                    judgment = claim.get('judgment', 'Unknown')
                    explanation = claim.get('explanation', 'No explanation provided')
                    
                    if judgment == 'Correct':
                        st.success(f"**Claim {i+1}:** {claim_text}")
                        st.write(f"*✅ {judgment}: {explanation}*")
                    elif judgment == 'Incorrect':
                        st.error(f"**Claim {i+1}:** {claim_text}")
                        st.write(f"*❌ {judgment}: {explanation}*")
                    else:  # Unsupported
                        st.warning(f"**Claim {i+1}:** {claim_text}")
                        st.write(f"*⚠️ {judgment}: {explanation}*")
                    st.markdown("---")
    else:
        st.info("No detailed correctness analysis available for claims distribution.")
    
    # Engagement metrics table (Requirement 3)
    st.markdown("**3. Engagement Metrics Table**")
    engagement = analysis_details.get('engagement', {})
    
    # Handle different engagement analysis structures
    full_analysis = None
    if 'full_analysis' in engagement:
        # When engagement agent succeeds, full_analysis is directly in engagement
        full_analysis = engagement.get('full_analysis', {})
    elif 'agent_analysis' in engagement:
        # Legacy structure where agent_analysis contains full_analysis
        agent_analysis = engagement.get('agent_analysis', {})
        full_analysis = agent_analysis.get('full_analysis', {})
    
    if full_analysis:
        quantitative = full_analysis.get('quantitative_metrics', {})
        qualitative = full_analysis.get('qualitative_analysis', {})
        
        # Create engagement metrics table structure as requested
        engagement_metrics = {
            "Quantity": {
                "StudentTalkRatio": f"{quantitative.get('student_talk_ratio', 0):.1f}%",
                "StudentTurnCount": quantitative.get('total_student_turns', 0),
                "AverageStudentTurnLength": f"{quantitative.get('average_turn_length', 0):.1f} words",
                "BackAndForthRatio": f"{quantitative.get('turns_per_10min', 0):.1f} per 10min"
            },
            "Quality": {
                "QuestionTypeDistribution": {
                    "Conceptual/Deep": qualitative.get('question_distribution', {}).get('conceptual_deep', 0),
                    "Clarification/Surface": qualitative.get('question_distribution', {}).get('clarification_surface', 0),
                    "Procedural/Admin": qualitative.get('question_distribution', {}).get('procedural_admin', 0)
                },
                "ElaborationIndex": qualitative.get('elaboration_index', 0),
                "DialogueDepth": qualitative.get('dialogue_depth', 0)
            },
            "Alignment": {
                "AverageTopicalOverlap": qualitative.get('topical_overlap', 0),
                "ContentCoverage": f"{qualitative.get('content_coverage', 0):.1f}%",
                "OffTopicRatio": f"{qualitative.get('off_topic_ratio', 0):.1f}%"
            },
            "Dynamics": {
                "EngagementDiversity": qualitative.get('engagement_diversity', 0),
                "TurnDistributionInequality": qualitative.get('turn_distribution_inequality', 0)
            }
        }
        
        # Display as nested tables
        for category, metrics in engagement_metrics.items():
            st.markdown(f"**{category} Metrics:**")
            if category == "Quality" and "QuestionTypeDistribution" in metrics:
                # Special handling for nested question distribution
                qtd = metrics["QuestionTypeDistribution"]
                quality_data = [
                    {"Metric": "Conceptual/Deep Questions", "Value": qtd["Conceptual/Deep"]},
                    {"Metric": "Clarification/Surface Questions", "Value": qtd["Clarification/Surface"]},
                    {"Metric": "Procedural/Admin Questions", "Value": qtd["Procedural/Admin"]},
                    {"Metric": "Elaboration Index", "Value": metrics["ElaborationIndex"]},
                    {"Metric": "Dialogue Depth", "Value": metrics["DialogueDepth"]}
                ]
                st.dataframe(pd.DataFrame(quality_data), hide_index=True, use_container_width=True)
            else:
                # Regular metrics display
                metric_data = [{"Metric": k, "Value": v} for k, v in metrics.items()]
                st.dataframe(pd.DataFrame(metric_data), hide_index=True, use_container_width=True)
            st.markdown("---")
    else:
        st.info("No detailed engagement metrics available.")
    
    # Topic time distribution (Requirement 4)
    st.markdown("**4. Lecture Coverage Analysis by Topic**")
    topic_coverage = analysis_details.get('topic_coverage', {})
    
    # Try to extract topic information from Topic Coverage analysis
    if topic_coverage and 'topic_analysis' in topic_coverage:
        topic_analysis = topic_coverage.get('topic_analysis', [])
        if topic_analysis:
            topic_data = []
            for topic_info in topic_analysis:
                topic_name = topic_info.get('topic', 'Unknown Topic')
                coverage_score = topic_info.get('coverage_score', 0)
                depth_score = topic_info.get('depth_score', 0)
                coverage_status = topic_info.get('coverage_status', 'unknown')
                
                # Convert status to more readable format
                status_display = {
                    'well_covered': '✅ Well Covered',
                    'partially_covered': '⚠️ Partially Covered',
                    'mentioned': '🔸 Mentioned',
                    'not_covered': '❌ Not Covered',
                    'unknown': '❓ Unknown'
                }.get(coverage_status, coverage_status)
                
                topic_data.append({
                    "Topic": topic_name,
                    "Status": status_display,
                    "Coverage Score": f"{coverage_score:.2f}",
                    "Depth Score": f"{depth_score:.2f}",
                    "Overall": f"{((coverage_score + depth_score) / 2):.2f}"
                })
            
            if topic_data:
                topic_df = pd.DataFrame(topic_data)
                st.dataframe(topic_df, hide_index=True, use_container_width=True)
                
                # Show overall summary
                st.markdown("**Summary:**")
                col1, col2, col3 = st.columns(3)

                overall_analysis = topic_coverage.get('overall_analysis', {})
                with col1:
                    st.metric("Average Coverage", f"{overall_analysis.get('coverage_score', 0):.2f}")
                with col2:
                    st.metric("Average Depth", f"{overall_analysis.get('depth_score', 0):.2f}")
                with col3:
                    analysis_method = topic_coverage.get('analysis_method', 'unknown')
                    method_display = "🤖 AI Analysis" if analysis_method == 'openai_enhanced' else "📊 Basic Analysis"
                    st.metric("Analysis Method", method_display)
            else:
                st.info("No topic analysis data available.")
        else:
            st.info("No topic analysis available - topics may not have been specified.")
    elif topic_coverage and 'scoring_details' in topic_coverage:
        # Show basic topic coverage information when detailed analysis isn't available
        scoring_details = topic_coverage.get('scoring_details', {})
        coverage_score = scoring_details.get('coverage_score', 0)
        depth_score = scoring_details.get('depth_score', 0)
        
        topic_summary_data = [
            {"Metric": "Overall Coverage Score", "Value": f"{coverage_score:.2f}"},
            {"Metric": "Content Depth Score", "Value": f"{depth_score:.2f}"},
            {"Metric": "Analysis Method", "Value": topic_coverage.get('analysis_method', 'fallback')}
        ]
        
        st.dataframe(pd.DataFrame(topic_summary_data), hide_index=True, use_container_width=True)
        st.info("💡 For detailed per-topic analysis, ensure topics are specified and OpenAI API is available.")
    else:
        st.info("Topic analysis requires specified topics and content evaluation.")
    
    # Recommendations
    if evaluation_report.get('recommendations'):
        st.subheader("💡 Recommendations")
        for rec in evaluation_report['recommendations']:
            if "⚠️" in rec or "❌" in rec:
                st.error(rec)
            elif "💡" in rec or "🎯" in rec:
                st.info(rec)
            elif "🌟" in rec or "👍" in rec:
                st.success(rec)
            else:
                st.info(rec)
    else:
        st.info("No specific recommendations generated.")


def show_analytics_dashboard():
    """Display analytics dashboard for principals"""
    st.header("📈 Analytics Dashboard")
    
    data_manager = get_data_manager()
    
    try:
        # Get school-wide analytics
        analytics = data_manager.get_school_analytics()
        
        if analytics and 'evaluations' in analytics:
            evaluations = analytics['evaluations']
            
            if evaluations:
                # Create DataFrame for analysis
                df = pd.DataFrame(evaluations)
                
                # Overview metrics
                col1, col2, col3, col4 = st.columns(4)
                
                with col1:
                    avg_score = df['evaluation_score'].mean()
                    st.metric("Average Score", f"{avg_score:.1f}", f"+{avg_score-75:.1f} vs target")
                
                with col2:
                    total_lectures = len(df)
                    st.metric("Total Lectures", total_lectures)
                
                with col3:
                    unique_teachers = df['teacher_name'].nunique()
                    st.metric("Active Teachers", unique_teachers)
                
                with col4:
                    recent_trend = df['evaluation_score'].tail(5).mean() - df['evaluation_score'].head(5).mean()
                    st.metric("Recent Trend", f"{recent_trend:+.1f}", "points")
                
                # Score distribution
                st.subheader("📊 Score Distribution")
                fig_hist = px.histogram(df, x='evaluation_score', nbins=20, 
                                       title="Distribution of Lecture Scores")
                st.plotly_chart(fig_hist, use_container_width=True)
                
                # Teacher performance comparison
                st.subheader("👨‍🏫 Teacher Performance")
                teacher_stats = df.groupby('teacher_name')['evaluation_score'].agg(['mean', 'count']).reset_index()
                teacher_stats.columns = ['Teacher', 'Average Score', 'Lectures Given']
                teacher_stats = teacher_stats.sort_values('Average Score', ascending=False)
                
                fig_bar = px.bar(teacher_stats, x='Teacher', y='Average Score', 
                               title="Average Scores by Teacher")
                st.plotly_chart(fig_bar, use_container_width=True)
                
                # Recent evaluations table
                st.subheader("📋 Recent Evaluations")
                recent_df = df.sort_values('evaluation_timestamp', ascending=False).head(10)
                display_cols = ['teacher_name', 'lecture_title', 'course_code', 'evaluation_score', 'evaluation_timestamp']
                st.dataframe(recent_df[display_cols], use_container_width=True)
                
            else:
                st.info("No evaluation data available yet. Upload some lectures to see analytics!")
        else:
            st.info("No analytics data available. Please ensure evaluations are being stored properly.")
            
    except Exception as e:
        st.error(f"Error loading analytics: {str(e)}")

def show_teacher_performance():
    """Display teacher performance analysis"""
    st.header("👨‍🏫 Teacher Performance Analysis")
    
    data_manager = get_data_manager()
    
    # Teacher selection
    try:
        analytics = data_manager.get_school_analytics()
        if analytics and 'evaluations' in analytics:
            df = pd.DataFrame(analytics['evaluations'])
            teachers = df['teacher_name'].unique().tolist()
            
            selected_teacher = st.selectbox("Select Teacher:", ["All Teachers"] + teachers, key="teacher_selectbox")
            
            if selected_teacher and selected_teacher != "All Teachers":
                # Get teacher-specific data
                teacher_evaluations = data_manager.get_evaluations_by_teacher(selected_teacher)
                
                if teacher_evaluations:
                    teacher_df = pd.DataFrame(teacher_evaluations)
                    
                    # Teacher overview
                    col1, col2, col3 = st.columns(3)
                    
                    with col1:
                        avg_score = teacher_df['evaluation_score'].mean()
                        st.metric("Average Score", f"{avg_score:.1f}")
                    
                    with col2:
                        total_lectures = len(teacher_df)
                        st.metric("Total Lectures", total_lectures)
                    
                    with col3:
                        improvement = teacher_df['evaluation_score'].iloc[-1] - teacher_df['evaluation_score'].iloc[0] if len(teacher_df) > 1 else 0
                        st.metric("Improvement", f"{improvement:+.1f}")
                    
                    # Performance over time
                    st.subheader("📈 Performance Trend")
                    fig_line = px.line(teacher_df, x='evaluation_timestamp', y='evaluation_score',
                                     title=f"{selected_teacher}'s Performance Over Time")
                    st.plotly_chart(fig_line, use_container_width=True)
                    
                    # Component analysis
                    st.subheader("📊 Component Analysis")
                    if 'score_breakdown' in teacher_df.columns:
                        # Extract component scores
                        components = ['Content Correctness', 'Class Engagement', 'Structure & Organization', 'Topic Coverage']
                        component_data = []
                        
                        for _, row in teacher_df.iterrows():
                            breakdown = row['score_breakdown']
                            if isinstance(breakdown, dict):
                                for comp in components:
                                    if comp in breakdown:
                                        component_data.append({
                                            'Component': comp,
                                            'Score': breakdown[comp],
                                            'Date': row['evaluation_timestamp']
                                        })
                        
                        if component_data:
                            comp_df = pd.DataFrame(component_data)
                            fig_comp = px.line(comp_df, x='Date', y='Score', color='Component',
                                             title="Component Scores Over Time")
                            st.plotly_chart(fig_comp, use_container_width=True)
                    
                    # Recent lectures
                    st.subheader("📋 Recent Lectures")
                    recent_lectures = teacher_df.sort_values('evaluation_timestamp', ascending=False).head(5)
                    display_cols = ['lecture_title', 'course_code', 'evaluation_score', 'evaluation_timestamp']
                    st.dataframe(recent_lectures[display_cols], use_container_width=True)
                    
                else:
                    st.info(f"No evaluation data found for {selected_teacher}")
            else:
                st.info("Please select a teacher to view detailed performance analysis.")
        else:
            st.info("No teacher data available yet.")
            
    except Exception as e:
        st.error(f"Error loading teacher data: {str(e)}")

# Main application
def main():
    """Main application function"""
    
    # Sidebar navigation
    st.sidebar.title("🎓 VirtuLearn Principal Dashboard")
    st.sidebar.markdown("---")
    
    page = st.sidebar.selectbox(
        "Navigate to:",
        ["📚 Lecture Upload", "📈 Analytics Dashboard", "👨‍🏫 Teacher Performance"],
        key="navigation_selectbox"
    )
    
    st.sidebar.markdown("---")
    st.sidebar.markdown("### System Status")
    
    # Check system status
    try:
        data_manager = get_data_manager()
        st.sidebar.success("✅ Database Connected")
    except Exception as e:
        st.sidebar.error("❌ Database Error")
        st.sidebar.error(str(e))
    
    # Check AI capabilities
    ai_status = []
    
    # Check OpenAI-powered fact checking
    try:
        from model.correctness_evaluator import OPENAI_AVAILABLE as fact_check_available
        if fact_check_available:
            ai_status.append("✅ AI-Powered Fact Checking")
        else:
            ai_status.append("⚠️ Fact Checking Fallback (No OpenAI API)")
    except ImportError:
        ai_status.append("⚠️ Fact Checking Fallback")
    
    # Check OpenAI-powered engagement analysis
    try:
        from model.engagement_evaluator import OPENAI_AVAILABLE as engagement_available
        if engagement_available:
            ai_status.append("✅ AI-Powered Engagement Analysis")
        else:
            ai_status.append("⚠️ Engagement Analysis Fallback (No OpenAI API)")
    except ImportError:
        ai_status.append("⚠️ Engagement Analysis Fallback")
    
    for status in ai_status:
        if "✅" in status:
            st.sidebar.success(status)
        else:
            st.sidebar.warning(status)
    
    st.sidebar.markdown("---")
    st.sidebar.markdown("*Powered by AI-Enhanced Evaluation*")
    
    # Main content area
    if page == "📚 Lecture Upload":
        show_lecture_upload()
    elif page == "📈 Analytics Dashboard":
        show_analytics_dashboard()
    elif page == "👨‍🏫 Teacher Performance":
        show_teacher_performance()


def create_topic_time_distribution_chart(evaluation_results):
    """Create a pie chart showing time distribution across topics"""
    structure_data = evaluation_results.get('structure', {})
    agent_analysis = structure_data.get('agent_analysis', {})
    full_analysis = agent_analysis.get('full_analysis', {})
    topic_breakdown = full_analysis.get('topic_breakdown', {})
    
    if not topic_breakdown:
        st.write("No topic breakdown data available.")
        return
    
    # Extract topics and estimated time distribution
    topics = []
    times = []
    
    for topic, details in topic_breakdown.items():
        if isinstance(details, dict):
            # Try to estimate time from coverage or assume equal distribution
            coverage = details.get('coverage_percentage', 0)
            if coverage > 0:
                topics.append(topic)
                times.append(coverage)
    
    if not topics:
        # Fallback: use topic names with equal distribution
        topics = list(topic_breakdown.keys())
        times = [100 / len(topics)] * len(topics)
    
    # Create pie chart
    fig = go.Figure(data=[go.Pie(
        labels=topics,
        values=times,
        hole=0.3,  # Donut chart
        textinfo='label+percent',
        textposition='auto',
        hovertemplate='<b>%{label}</b><br>Time: %{percent}<br><extra></extra>',
        marker=dict(
            colors=px.colors.qualitative.Set3[:len(topics)]
        )
    )])
    
    fig.update_layout(
        title="Lecture Time Distribution by Topic",
        font=dict(size=12),
        height=400,
        showlegend=True,
        legend=dict(
            orientation="v",
            yanchor="middle",
            y=0.5,
            xanchor="left",
            x=1.01
        )
    )
    
    st.plotly_chart(fig, use_container_width=True)


def create_enhanced_feedback_analysis(evaluation_results):
    """Create detailed feedback analysis for each evaluation category"""
    categories = {
        'Topic Coverage': evaluation_results.get('topic_coverage', {}),
        'Engagement': evaluation_results.get('engagement', {}),
        'Structure': evaluation_results.get('structure', {}),
        'Correctness': evaluation_results.get('correctness', {})
    }
    
    st.subheader("🔍 Detailed Category Analysis")
    
    # Create tabs for each category
    tabs = st.tabs([f"📊 {cat}" for cat in categories.keys()])
    
    for i, (category, data) in enumerate(categories.items()):
        with tabs[i]:
            agent_analysis = data.get('agent_analysis', {})
            full_analysis = agent_analysis.get('full_analysis', {})
            
            # Display score prominently - using actual score from evaluation
            score = 0
            if category == 'Topic Coverage':
                score = agent_analysis.get('topic_coverage_score', 0)
            elif category == 'Engagement':
                score = agent_analysis.get('engagement_score', 0)
            elif category == 'Structure':
                score = agent_analysis.get('structure_score', 0)
            elif category == 'Correctness':
                score = agent_analysis.get('correctness_score', 0)
            
            col1, col2 = st.columns([1, 3])
            
            with col1:
                # Score gauge
                fig = go.Figure(go.Indicator(
                    mode = "gauge+number+delta",
                    value = score,
                    domain = {'x': [0, 1], 'y': [0, 1]},
                    title = {'text': f"{category} Score"},
                    gauge = {
                        'axis': {'range': [None, 100]},
                        'bar': {'color': "darkblue"},
                        'steps': [
                            {'range': [0, 50], 'color': "lightgray"},
                            {'range': [50, 80], 'color': "yellow"},
                            {'range': [80, 100], 'color': "green"}],
                        'threshold': {
                            'line': {'color': "red", 'width': 4},
                            'thickness': 0.75,
                            'value': 90}}))
                fig.update_layout(height=200, margin=dict(l=20, r=20, t=40, b=20))
                st.plotly_chart(fig, use_container_width=True)
            
            with col2:
                # Category-specific insights
                if category == 'Topic Coverage':
                    st.write("**Key Areas:**")
                    clarity = full_analysis.get('clarity_assessment', {})
                    depth = full_analysis.get('depth_analysis', {})
                    accuracy = full_analysis.get('accuracy_check', {})
                    
                    if clarity:
                        st.write(f"• **Clarity:** {clarity.get('score', 'N/A')}/10 - {clarity.get('feedback', 'No feedback available')}")
                    if depth:
                        st.write(f"• **Depth:** {depth.get('score', 'N/A')}/10 - {depth.get('feedback', 'No feedback available')}")
                    if accuracy:
                        st.write(f"• **Accuracy:** {accuracy.get('score', 'N/A')}/10 - {accuracy.get('feedback', 'No feedback available')}")
                
                elif category == 'Engagement':
                    st.write("**Engagement Metrics:**")
                    quantitative = full_analysis.get('quantitative_metrics', {})
                    qualitative = full_analysis.get('qualitative_assessment', {})
                    
                    if quantitative:
                        talk_ratio = quantitative.get('student_talk_ratio', 0)
                        inferred_q = quantitative.get('inferred_student_questions', quantitative.get('total_student_turns', 0))
                        
                        st.write(f"• **Student Talk Ratio:** {talk_ratio:.1f}% {'(inferred from teacher responses)' if talk_ratio > 0 else '(no questions detected)'}")
                        st.write(f"• **Inferred Questions:** {inferred_q} questions detected")
                        st.write(f"• **Questions per 10min:** {quantitative.get('turns_per_10min', 0):.1f}")
                        
                        # Show teacher response indicators if available
                        response_indicators = quantitative.get('teacher_response_indicators', 0)
                        if response_indicators > 0:
                            st.write(f"• **Teacher Response Cues:** {response_indicators} detected")
                    
                    if qualitative:
                        st.write(f"• **Interaction Quality:** {qualitative.get('interaction_quality', 'Not assessed')}")
                
                elif category == 'Structure':
                    st.write("**Structure Elements:**")
                    flow = full_analysis.get('flow_analysis', {})
                    organization = full_analysis.get('organization_score', {})
                    transitions = full_analysis.get('transition_quality', {})
                    
                    if flow:
                        st.write(f"• **Flow:** {flow.get('score', 'N/A')}/10 - {flow.get('feedback', 'No feedback available')}")
                    if organization:
                        st.write(f"• **Organization:** {organization.get('score', 'N/A')}/10")
                    if transitions:
                        st.write(f"• **Transitions:** {transitions.get('score', 'N/A')}/10")
                
                elif category == 'Correctness':
                    st.write("**Fact Checking Results:**")
                    fact_checks = full_analysis.get('fact_checks', [])
                    
                    if fact_checks:
                        correct_count = sum(1 for fc in fact_checks if fc.get('verdict') == 'CORRECT')
                        total_count = len(fact_checks)
                        st.write(f"• **Verified Claims:** {correct_count}/{total_count}")
                        
                        # Show problematic claims
                        incorrect = [fc for fc in fact_checks if fc.get('verdict') in ['INCORRECT', 'QUESTIONABLE']]
                        if incorrect:
                            st.write("• **Issues Found:**")
                            for fc in incorrect[:3]:  # Show first 3
                                st.write(f"  - {fc.get('claim', 'Unknown claim')[:100]}...")
            
            # Detailed recommendations with transcript evidence
            st.write("**🎯 Specific Recommendations:**")
            recommendations = full_analysis.get('recommendations', [])
            if recommendations:
                for i, rec in enumerate(recommendations[:5], 1):  # Show top 5
                    st.write(f"{i}. {rec}")
            else:
                # Generate category-specific default recommendations with evidence
                if category == 'Topic Coverage':
                    topic_analysis = full_analysis.get('topics_analysis', {})
                    if topic_analysis.get('uncovered_topics'):
                        st.write(f"1. Address missing topics: {', '.join(topic_analysis['uncovered_topics'][:2])}")
                    else:
                        st.write("1. Consider adding more examples to illustrate complex concepts")
                    st.write("2. Ensure definitions are clear and accessible to your audience")
                    st.write("3. Use analogies to make abstract concepts more concrete")
                elif category == 'Engagement':
                    quantitative = full_analysis.get('quantitative_metrics', {})
                    student_talk = quantitative.get('student_talk_ratio', 0)
                    inferred_q = quantitative.get('inferred_student_questions', quantitative.get('total_student_turns', 0))
                    
                    if student_talk < 3:
                        st.write(f"1. No student questions detected - encourage participation through direct prompts")
                    elif student_talk < 8:
                        st.write(f"1. Limited engagement ({student_talk:.1f}%) - create more question opportunities")
                    else:
                        st.write(f"1. Good engagement level ({student_talk:.1f}%) - maintain current interaction style")
                    
                    if inferred_q < 2:
                        st.write("2. Use more explicit question prompts: 'Any questions?' or 'What do you think?'")
                    else:
                        st.write("2. Continue encouraging questions - students are actively participating")
                    
                    st.write("3. Look for raised hands and pause for questions after complex topics")
                elif category == 'Structure':
                    flow_analysis = full_analysis.get('flow_analysis', {})
                    if flow_analysis.get('score', 10) < 7:
                        st.write(f"1. Improve transitions: {flow_analysis.get('feedback', 'Add clearer connections between topics')}")
                    else:
                        st.write("1. Provide clear agenda at the beginning of the lecture")
                    st.write("2. Use signposting to help students follow the flow")
                    st.write("3. Include regular summaries of key points")
                elif category == 'Correctness':
                    fact_checks = full_analysis.get('fact_checks', [])
                    incorrect_count = len([fc for fc in fact_checks if fc.get('verdict') in ['INCORRECT', 'QUESTIONABLE']])
                    if incorrect_count > 0:
                        st.write(f"1. Review {incorrect_count} flagged claims for accuracy")
                    else:
                        st.write("1. Double-check statistical claims and data")
                    st.write("2. Cite sources for factual statements")
                    st.write("3. Consider peer review for technical content")


if __name__ == "__main__":
    main()